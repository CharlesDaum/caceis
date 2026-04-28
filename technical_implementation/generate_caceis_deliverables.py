from __future__ import annotations

from pathlib import Path
import math
import textwrap

import matplotlib.pyplot as plt
import nbformat as nbf
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT_DIR = ROOT / "deliverables"
FIG_DIR = OUTPUT_DIR / "figures"
TABLE_DIR = OUTPUT_DIR / "tables"


def ensure_dirs() -> None:
    OUTPUT_DIR.mkdir(exist_ok=True)
    FIG_DIR.mkdir(exist_ok=True)
    TABLE_DIR.mkdir(exist_ok=True)


def clean_columns(df: pd.DataFrame) -> pd.DataFrame:
    out = df.copy()
    out.columns = [str(c).strip().replace("\xa0", " ") for c in out.columns]
    return out


def load_sources() -> dict[str, pd.DataFrame]:
    workforce = clean_columns(pd.read_excel(ROOT / "HR Data" / "Data.xlsx", sheet_name="Sheet1"))
    performance_2023 = clean_columns(
        pd.read_excel(ROOT / "HR Data" / "20240222 - CACEIS Notes evaluation 2023.xlsx")
    )
    performance_2024 = clean_columns(
        pd.read_excel(
            ROOT / "HR Data" / "20250218 - Stats CACEIS EAE EP 18-02-2025 Version Définitive cloture.xlsx",
            sheet_name="Database",
        )
    )
    absenteeism_2025 = clean_columns(
        pd.read_excel(
            ROOT / "HR Data" / "20260121 - Absentéisme_-_détail_affectation_-_Bilan_social 2025.xlsx",
            sheet_name="extract",
        )
    )
    training_records = clean_columns(pd.read_excel(ROOT / "Training" / "Training_Records_Unnamed.xlsx"))
    training_quick = clean_columns(pd.read_excel(ROOT / "Training" / "Quick_Review_Unnamed.xlsx"))
    training_cold = clean_columns(pd.read_excel(ROOT / "Training" / "Cold_Review_Unnamed.xlsx"))
    finance_pl = clean_columns(pd.read_excel(ROOT / "Finance" / "AlbertSchool_CACEIS_PL-FTE_22-25_Sent.xlsx", sheet_name="Synthese_PL"))
    finance_fte = clean_columns(pd.read_excel(ROOT / "Finance" / "AlbertSchool_CACEIS_PL-FTE_22-25_Sent.xlsx", sheet_name="Synthese_ETP"))
    return {
        "workforce": workforce,
        "performance_2023": performance_2023,
        "performance_2024": performance_2024,
        "absenteeism_2025": absenteeism_2025,
        "training_records": training_records,
        "training_quick": training_quick,
        "training_cold": training_cold,
        "finance_pl": finance_pl,
        "finance_fte": finance_fte,
    }


def prepare_workforce(df: pd.DataFrame) -> pd.DataFrame:
    out = df.copy()
    out["PERIOD"] = pd.to_datetime(out["PERIOD"], errors="coerce")
    out["employee_id"] = out["ID Employee"].fillna(out["ID Employee.1"])
    out["country"] = out["COUNTRY_GROUP_LABEL_EN"]
    out["gender"] = out["SEXE_GROUP_LABEL_EN"]
    out["contract"] = out["CONTRACT_GROUP_LABEL_EN"]
    out["entity"] = out["ENTITY_LABEL_LOCAL"]
    out["age_range"] = out["Age range"]
    out = out.dropna(subset=["employee_id", "PERIOD"])
    return out


def prepare_performance_2023(df: pd.DataFrame) -> pd.DataFrame:
    out = df.copy()
    out["country"] = out["Pays"].ffill()
    out["year"] = 2023
    out["employee_id"] = out["IUG"]
    out["performance_score"] = pd.to_numeric(out["Note"], errors="coerce")
    return out[["employee_id", "country", "year", "performance_score"]].dropna(subset=["employee_id"])


def prepare_performance_2024(df: pd.DataFrame) -> pd.DataFrame:
    out = df.copy()
    out["employee_id"] = out["IUG"]
    out["year"] = pd.to_numeric(out["Année"], errors="coerce")
    out["score"] = pd.to_numeric(out["Note de performance"], errors="coerce")
    out["country"] = out["Country"]
    doc = out["Nom du document"].fillna("")
    perf_mask = (
        doc.str.contains("Evaluation", case=False, na=False)
        | doc.str.contains("Responsabilit", case=False, na=False)
        | out["score"].notna()
    )
    perf_docs = out[perf_mask].copy()
    by_employee = (
        perf_docs.sort_values(["employee_id", "score"], ascending=[True, False])
        .groupby(["employee_id", "country", "year"], dropna=False, as_index=False)
        .agg(
            performance_score=("score", "max"),
            review_completion_status=("Statut", lambda s: s.dropna().iloc[0] if not s.dropna().empty else None),
            document_status=("Statut du document", lambda s: s.dropna().iloc[0] if not s.dropna().empty else None),
        )
    )
    return by_employee


def prepare_absenteeism(df: pd.DataFrame) -> pd.DataFrame:
    out = df.copy()
    out["Date Absence"] = pd.to_datetime(out["Date Absence"], errors="coerce")
    out["employee_id"] = out["Employee Code"]
    out["calendar_days"] = pd.to_numeric(out["Jour Calendaires Absence"], errors="coerce").fillna(0)
    out["working_days"] = pd.to_numeric(out["Jours Ouvrés Absence"], errors="coerce").fillna(0)
    out["absence_group"] = out["Regroupement Jour Absences"].fillna("Unknown")
    out["absence_reason"] = out["Motif Jour Absence"].fillna("Unknown")
    out["entity"] = out["Société"]
    out["gender"] = out["Genre"]
    return out.dropna(subset=["employee_id", "Date Absence"])


def prepare_training_records(df: pd.DataFrame) -> pd.DataFrame:
    out = df.copy()
    out["employee_id"] = out["Employee Code"]
    out["year"] = pd.to_numeric(out["Year"], errors="coerce")
    out["hours"] = pd.to_numeric(out["Total_Training_Hours"], errors="coerce").fillna(0)
    out["status"] = out["Status"].fillna("Unknown")
    out["entity"] = out["Entity"]
    out["certification_flag"] = out["Certifications"].astype(str).str.lower().eq("yes")
    return out


def prepare_training_quick(df: pd.DataFrame) -> pd.DataFrame:
    out = df.copy()
    out["employee_id"] = out["Matricule"]
    out["general_rating"] = pd.to_numeric(out["Note générale"], errors="coerce")
    out["status"] = out["Statut"]
    return out


def prepare_training_cold(df: pd.DataFrame) -> pd.DataFrame:
    out = df.copy()
    out["employee_id"] = out["Matricule"]
    out["status"] = out["Status"]
    impact_cols = [
        "Considérez-vous que cette formation vous a permis de prendre confiance en vous :",
        "Considérez-vous que cette formation vous a permis de faciliter votre quotidien :",
        "Considérez-vous que cette formation vous a permis d’améliorer la qualité ou l’efficacité de votre travail :",
        "Considérez-vous que cette formation vous a permis de vous perfectionner dans un domaine que vous connaissiez déjà :",
        "Considérez-vous que cette formation vous a permis de développer de nouvelles compétences :",
    ]
    normalized = {}
    for col in out.columns:
        normalized[col.replace("\xa0", " ")] = col
    selected = [normalized[c] for c in impact_cols if c in normalized]
    out["positive_impact_answers"] = out[selected].apply(
        lambda row: sum(str(v).strip().lower().startswith("oui") for v in row if pd.notna(v)), axis=1
    )
    out["impact_answer_count"] = out[selected].apply(lambda row: sum(pd.notna(v) for v in row), axis=1)
    return out


def extract_finance_metrics(pl_df: pd.DataFrame, fte_df: pd.DataFrame) -> pd.DataFrame:
    pl = pl_df.iloc[:, :5].copy()
    pl.columns = ["metric", "2022", "2023", "2024", "2025"]
    pl = pl[pl["metric"].notna()]
    pl = pl.melt(id_vars="metric", var_name="year", value_name="value")

    fte = fte_df.iloc[1:4, [0, 2, 4, 6, 8]].copy()
    fte.columns = ["segment", "2022", "2023", "2024", "2025"]
    fte["segment"] = ["total", "permanent", "temporary"]
    fte = fte.melt(id_vars="segment", var_name="year", value_name="value")

    total_fte = fte[fte["segment"] == "total"][["year", "value"]].rename(columns={"value": "avg_fte_proxy"})
    nbi = pl[pl["metric"] == "Net Banking Income (PNB)"][["year", "value"]].rename(columns={"value": "nbi"})
    personnel = pl[pl["metric"] == "Rémunérations & charges"][["year", "value"]].rename(columns={"value": "personnel_cost"})
    training_cost = pl[pl["metric"] == "Formation (training costs)"][["year", "value"]].rename(columns={"value": "training_cost"})

    finance = total_fte.merge(nbi, on="year").merge(personnel, on="year").merge(training_cost, on="year")
    finance["nbi_per_fte"] = finance["nbi"] / finance["avg_fte_proxy"]
    finance["personnel_cost_per_fte"] = finance["personnel_cost"].abs() / finance["avg_fte_proxy"]
    finance["training_cost_per_fte"] = finance["training_cost"].abs() / finance["avg_fte_proxy"]
    finance["year"] = finance["year"].astype(int)
    return finance.sort_values("year")


def build_kpis(data: dict[str, pd.DataFrame]) -> tuple[pd.DataFrame, dict[str, pd.DataFrame]]:
    workforce = prepare_workforce(data["workforce"])
    performance_2023 = prepare_performance_2023(data["performance_2023"])
    performance_2024 = prepare_performance_2024(data["performance_2024"])
    absenteeism = prepare_absenteeism(data["absenteeism_2025"])
    training_records = prepare_training_records(data["training_records"])
    training_quick = prepare_training_quick(data["training_quick"])
    training_cold = prepare_training_cold(data["training_cold"])
    finance = extract_finance_metrics(data["finance_pl"], data["finance_fte"])

    monthly_headcount = (
        workforce.groupby("PERIOD", as_index=False)
        .agg(headcount=("employee_id", "nunique"), female_share=("gender", lambda s: (s == "Female").mean()))
        .sort_values("PERIOD")
    )
    latest_workforce = workforce[workforce["PERIOD"] == workforce["PERIOD"].max()].copy()
    headcount_by_country = (
        latest_workforce.groupby("country", as_index=False)
        .agg(headcount=("employee_id", "nunique"), female_share=("gender", lambda s: (s == "Female").mean()))
        .sort_values("headcount", ascending=False)
    )

    reviews_2024 = performance_2024.copy()
    reviews_2024["completed_flag"] = reviews_2024["review_completion_status"].eq("Terminé")
    review_metrics = (
        reviews_2024.groupby("country", as_index=False)
        .agg(
            employees_with_reviews=("employee_id", "nunique"),
            completion_rate=("completed_flag", "mean"),
            avg_performance_score=("performance_score", "mean"),
            high_performer_share=("performance_score", lambda s: s.ge(4).mean()),
        )
        .sort_values("employees_with_reviews", ascending=False)
    )

    training_done = training_records[training_records["status"] == "Réalisé"].copy()
    training_kpis = (
        training_done.groupby("year", as_index=False)
        .agg(
            trained_employees=("employee_id", "nunique"),
            completed_courses=("Attended_Courses", "count"),
            training_hours=("hours", "sum"),
            avg_hours_per_trained_employee=("hours", lambda s: s.sum() / max(training_done.loc[s.index, "employee_id"].nunique(), 1)),
            certification_rate=("certification_flag", "mean"),
        )
        .sort_values("year")
    )

    quick_completed = training_quick[training_quick["status"] == "Complétée"]
    quick_scores = (
        quick_completed.groupby(pd.to_datetime(quick_completed["Date"], dayfirst=True, errors="coerce").dt.year, as_index=False)
        .agg(avg_training_rating=("general_rating", "mean"), rating_responses=("general_rating", "count"))
        .rename(columns={"Date": "year"})
        .sort_values("year")
    )

    cold_completed = training_cold[training_cold["status"] == "Complétée"].copy()
    cold_completed["impact_positive_rate"] = cold_completed["positive_impact_answers"] / cold_completed["impact_answer_count"].replace(0, pd.NA)
    cold_impact = (
        cold_completed.groupby(pd.to_datetime(cold_completed["Date"], dayfirst=True, errors="coerce").dt.year, as_index=False)
        .agg(avg_positive_impact_rate=("impact_positive_rate", "mean"), impact_responses=("impact_answer_count", "sum"))
        .rename(columns={"Date": "year"})
        .sort_values("year")
    )

    workforce_2025_avg = monthly_headcount[monthly_headcount["PERIOD"].dt.year == 2025]["headcount"].mean()
    latest_headcount = monthly_headcount.iloc[-1]["headcount"]
    absence_by_employee = (
        absenteeism.groupby("employee_id", as_index=False)
        .agg(
            absence_days=("working_days", "sum"),
            absence_events=("Date Absence", "count"),
            main_group=("absence_group", lambda s: s.value_counts().idxmax()),
        )
    )
    absence_group_mix = (
        absenteeism.groupby("absence_group", as_index=False)
        .agg(absence_days=("working_days", "sum"))
        .sort_values("absence_days", ascending=False)
    )

    cross_training_perf = training_done.groupby("employee_id", as_index=False).agg(training_hours=("hours", "sum"))
    performance_2024_scores = reviews_2024[["employee_id", "country", "performance_score"]].dropna(subset=["performance_score"])
    training_perf = performance_2024_scores.merge(cross_training_perf, on="employee_id", how="left")
    training_perf["training_hours"] = training_perf["training_hours"].fillna(0)
    training_perf["training_band"] = pd.cut(
        training_perf["training_hours"],
        bins=[-0.1, 0.1, 7, 20, math.inf],
        labels=["0h", "0.1-7h", "7-20h", "20h+"],
    )
    training_perf_summary = (
        training_perf.groupby("training_band", observed=False, as_index=False)
        .agg(avg_performance_score=("performance_score", "mean"), employee_count=("employee_id", "nunique"))
    )

    kpi_rows = [
        {"kpi": "Latest monthly headcount", "value": int(latest_headcount), "definition": "Unique active employees in latest workforce snapshot"},
        {"kpi": "Female representation", "value": round(float(monthly_headcount.iloc[-1]["female_share"]) * 100, 1), "definition": "Share of active workforce labeled Female in latest snapshot (%)"},
        {"kpi": "2024 review completion rate", "value": round(float(reviews_2024["completed_flag"].mean()) * 100, 1), "definition": "Share of employees with a completed 2024 performance review (%)"},
        {"kpi": "2024 average performance score", "value": round(float(reviews_2024["performance_score"].mean()), 2), "definition": "Mean available 2024 performance rating (1-5)"},
        {"kpi": "2025 training hours", "value": round(float(training_done[training_done["year"] == 2025]["hours"].sum()), 1), "definition": "Total completed training hours in 2025"},
        {"kpi": "Training satisfaction", "value": round(float(quick_completed["general_rating"].mean()), 2), "definition": "Average post-training rating across completed quick reviews (1-5)"},
        {"kpi": "2025 absence days per active employee", "value": round(float(absence_by_employee["absence_days"].sum() / latest_headcount), 2), "definition": "Total 2025 worked absence days divided by latest active headcount"},
        {"kpi": "2025 NBI per FTE", "value": round(float(finance.loc[finance["year"] == 2025, "nbi_per_fte"].iloc[0]), 2), "definition": "Net banking income divided by average FTE"},
    ]
    kpi_table = pd.DataFrame(kpi_rows)

    supporting = {
        "monthly_headcount": monthly_headcount,
        "headcount_by_country": headcount_by_country,
        "review_metrics": review_metrics,
        "training_kpis": training_kpis,
        "quick_scores": quick_scores,
        "cold_impact": cold_impact,
        "absence_by_employee": absence_by_employee,
        "absence_group_mix": absence_group_mix,
        "finance": finance,
        "training_perf_summary": training_perf_summary,
        "performance_2023": performance_2023,
        "performance_2024": performance_2024_scores,
    }
    return kpi_table, supporting


def save_tables(kpi_table: pd.DataFrame, supporting: dict[str, pd.DataFrame]) -> None:
    kpi_table.to_csv(TABLE_DIR / "kpi_summary.csv", index=False)
    for name, df in supporting.items():
        df.to_csv(TABLE_DIR / f"{name}.csv", index=False)


def style_axis(ax: plt.Axes) -> None:
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.grid(axis="y", alpha=0.2)


def save_figures(supporting: dict[str, pd.DataFrame]) -> None:
    palette = {
        "navy": "#16324F",
        "teal": "#1F7A8C",
        "sand": "#F2D0A9",
        "coral": "#E07A5F",
        "slate": "#4F6D7A",
    }

    headcount = supporting["monthly_headcount"]
    fig, ax1 = plt.subplots(figsize=(10, 5))
    ax1.plot(headcount["PERIOD"], headcount["headcount"], color=palette["navy"], linewidth=2.5)
    ax1.set_title("Monthly Workforce Capacity")
    ax1.set_ylabel("Headcount")
    style_axis(ax1)
    ax2 = ax1.twinx()
    ax2.plot(headcount["PERIOD"], headcount["female_share"] * 100, color=palette["coral"], linewidth=2)
    ax2.set_ylabel("Female share (%)")
    fig.tight_layout()
    fig.savefig(FIG_DIR / "headcount_and_gender.png", dpi=180)
    plt.close(fig)

    reviews = supporting["review_metrics"]
    fig, ax = plt.subplots(figsize=(8, 5))
    plot_df = reviews.sort_values("completion_rate", ascending=True)
    ax.barh(plot_df["country"], plot_df["completion_rate"] * 100, color=palette["teal"])
    ax.set_title("2024 Review Completion by Country")
    ax.set_xlabel("Completion rate (%)")
    style_axis(ax)
    fig.tight_layout()
    fig.savefig(FIG_DIR / "review_completion_by_country.png", dpi=180)
    plt.close(fig)

    training = supporting["training_kpis"]
    fig, ax = plt.subplots(figsize=(8, 5))
    ax.bar(training["year"].astype(str), training["training_hours"], color=palette["sand"])
    ax.set_title("Completed Training Hours")
    ax.set_xlabel("Year")
    ax.set_ylabel("Hours")
    style_axis(ax)
    fig.tight_layout()
    fig.savefig(FIG_DIR / "training_hours.png", dpi=180)
    plt.close(fig)

    satisfaction = supporting["quick_scores"].merge(supporting["cold_impact"], on="year", how="outer").sort_values("year")
    fig, ax1 = plt.subplots(figsize=(8, 5))
    ax1.plot(satisfaction["year"], satisfaction["avg_training_rating"], marker="o", color=palette["navy"])
    ax1.set_ylabel("Quick review rating (1-5)")
    ax1.set_xlabel("Year")
    ax1.set_title("Training Effectiveness Signals")
    ax2 = ax1.twinx()
    ax2.plot(satisfaction["year"], satisfaction["avg_positive_impact_rate"] * 100, marker="s", color=palette["coral"])
    ax2.set_ylabel("Positive impact response (%)")
    style_axis(ax1)
    fig.tight_layout()
    fig.savefig(FIG_DIR / "training_effectiveness.png", dpi=180)
    plt.close(fig)

    absence = supporting["absence_group_mix"].head(8)
    fig, ax = plt.subplots(figsize=(8, 5))
    ax.barh(absence["absence_group"], absence["absence_days"], color=palette["slate"])
    ax.set_title("2025 Absence Day Mix")
    ax.set_xlabel("Worked absence days")
    style_axis(ax)
    fig.tight_layout()
    fig.savefig(FIG_DIR / "absence_mix.png", dpi=180)
    plt.close(fig)

    finance = supporting["finance"]
    fig, ax = plt.subplots(figsize=(8, 5))
    ax.plot(finance["year"], finance["nbi_per_fte"], marker="o", color=palette["navy"], label="NBI per FTE")
    ax.plot(finance["year"], finance["personnel_cost_per_fte"], marker="o", color=palette["coral"], label="Personnel cost per FTE")
    ax.set_title("Value Creation Efficiency")
    ax.set_xlabel("Year")
    ax.set_ylabel("Amount")
    ax.legend(frameon=False)
    style_axis(ax)
    fig.tight_layout()
    fig.savefig(FIG_DIR / "value_creation_efficiency.png", dpi=180)
    plt.close(fig)

    perf = supporting["training_perf_summary"]
    fig, ax = plt.subplots(figsize=(8, 5))
    ax.bar(perf["training_band"].astype(str), perf["avg_performance_score"], color=palette["teal"])
    ax.set_title("Performance Score by Training Intensity")
    ax.set_xlabel("Training hours band")
    ax.set_ylabel("Average performance score")
    style_axis(ax)
    fig.tight_layout()
    fig.savefig(FIG_DIR / "training_vs_performance.png", dpi=180)
    plt.close(fig)


def make_notebook() -> None:
    nb = nbf.v4.new_notebook()
    cells = []
    cells.append(
        nbf.v4.new_markdown_cell(
            "# CACEIS Human Capital Valuation\n"
            "This notebook documents the reproducible pipeline used to transform the CACEIS HR datasets into a human capital valuation framework."
        )
    )
    cells.append(
        nbf.v4.new_markdown_cell(
            "## Scope\n"
            "- Source systems: workforce snapshots, performance reviews, training, absenteeism, finance\n"
            "- Output: KPI framework, exploratory analysis, preliminary AI prototype direction"
        )
    )
    cells.append(
        nbf.v4.new_code_cell(
            "from pathlib import Path\n"
            "import pandas as pd\n"
            "cwd = Path.cwd().resolve()\n"
            "if (cwd / 'deliverables' / 'tables').exists():\n"
            "    ROOT = cwd\n"
            "elif (cwd / 'tables').exists() and cwd.name == 'deliverables':\n"
            "    ROOT = cwd.parent\n"
            "else:\n"
            "    ROOT = cwd\n"
            "TABLE_DIR = ROOT / 'deliverables' / 'tables'\n"
            "FIG_DIR = ROOT / 'deliverables' / 'figures'\n"
            "kpis = pd.read_csv(TABLE_DIR / 'kpi_summary.csv')\n"
            "kpis"
        )
    )
    cells.append(
        nbf.v4.new_markdown_cell(
            "## KPI logic\n"
            "1. Workforce capacity and inclusion: headcount trend and female representation.\n"
            "2. Performance governance: completion rate and score distribution for annual reviews.\n"
            "3. Capability building: completed training hours, satisfaction, and perceived impact.\n"
            "4. Well-being risk: worked absence days per employee and absence mix.\n"
            "5. Value creation efficiency: net banking income and personnel cost per FTE."
        )
    )
    for name in [
        "monthly_headcount",
        "review_metrics",
        "training_kpis",
        "quick_scores",
        "cold_impact",
        "absence_group_mix",
        "finance",
        "training_perf_summary",
    ]:
        cells.append(
            nbf.v4.new_code_cell(
                f"pd.read_csv(TABLE_DIR / '{name}.csv').head(10)"
            )
        )
    cells.append(
        nbf.v4.new_markdown_cell(
            "## Generated charts\n"
            "The pipeline exports slide-ready charts to `deliverables/figures/`."
        )
    )
    for fig_name in [
        "headcount_and_gender.png",
        "review_completion_by_country.png",
        "training_hours.png",
        "training_effectiveness.png",
        "absence_mix.png",
        "value_creation_efficiency.png",
        "training_vs_performance.png",
    ]:
        cells.append(
            nbf.v4.new_code_cell(
                "from IPython.display import Image, display\n"
                f"display(Image(filename=str(FIG_DIR / '{fig_name}')))"
            )
        )
    cells.append(
        nbf.v4.new_markdown_cell(
            "## Governance and AI readiness\n"
            "- Privacy: analysis stays at anonymized employee ID level.\n"
            "- Bias controls: country and gender segment checks are retained for fairness diagnostics.\n"
            "- Prototype direction: rule-based KPI monitoring plus predictive risk scoring for absenteeism/performance coverage gaps."
        )
    )
    nb["cells"] = cells
    nb["metadata"]["kernelspec"] = {"display_name": "Python 3", "language": "python", "name": "python3"}
    nb["metadata"]["language_info"] = {"name": "python", "version": "3.x"}
    with open(OUTPUT_DIR / "CACEIS_Human_Capital_Implementation.ipynb", "w", encoding="utf-8") as f:
        nbf.write(nb, f)


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(11.8), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Aptos Display"
    run.font.bold = True
    run.font.size = Pt(26)
    run.font.color.rgb = RGBColor(22, 50, 79)
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(11.8), Inches(0.5))
        p2 = sub_box.text_frame.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Aptos"
        r2.font.size = Pt(11)
        r2.font.color.rgb = RGBColor(79, 109, 122)


def add_bullets(slide, items: list[str], left: float, top: float, width: float, height: float, font_size: int = 18) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(30, 30, 30)
        p.space_after = Pt(8)


def add_image(slide, image_path: Path, left: float, top: float, width: float) -> None:
    slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width))


def add_kpi_cards(slide, kpis: pd.DataFrame) -> None:
    colors = [
        RGBColor(22, 50, 79),
        RGBColor(31, 122, 140),
        RGBColor(224, 122, 95),
        RGBColor(79, 109, 122),
    ]
    for i, (_, row) in enumerate(kpis.head(4).iterrows()):
        left = 0.6 + i * 3.0
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.5), Inches(2.6), Inches(1.4))
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = colors[i]
        shape.line.color.rgb = colors[i]
        tf = shape.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = str(row["value"])
        r1.font.name = "Aptos Display"
        r1.font.bold = True
        r1.font.size = Pt(22)
        r1.font.color.rgb = RGBColor(255, 255, 255)
        p1.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = row["kpi"]
        p2.font.name = "Aptos"
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(255, 255, 255)
        p2.alignment = PP_ALIGN.CENTER


def make_presentation(kpis: pd.DataFrame, supporting: dict[str, pd.DataFrame]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "CACEIS Human Capital Valuation Framework", "From raw HR records to an AI-ready valuation layer")
    add_kpi_cards(slide, kpis)
    add_bullets(
        slide,
        [
            "Objective: shift HR data from cost tracking toward capacity, capability, performance, well-being and value creation.",
            "Data foundation: workforce snapshots, performance reviews, training records, absenteeism logs, and P&L/FTE history.",
            "Deliverable logic: reproducible KPI pipeline + exploratory insights + AI prototype direction.",
        ],
        0.8,
        3.3,
        11.7,
        2.8,
        18,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Framework", "Seven KPIs grouped around business value, not only cost")
    add_bullets(
        slide,
        [
            "Workforce capacity: monthly headcount to monitor operating scale and resilience.",
            "Inclusion: female representation to track talent mix and fairness at workforce level.",
            "Performance governance: annual review completion rate to measure managerial coverage.",
            "Performance quality: mean rating and high-performer share to capture contribution quality.",
            "Capability investment: completed training hours and certification incidence.",
            "Learning effectiveness: training satisfaction and perceived on-the-job impact.",
            "Well-being and efficiency: absence days per employee plus NBI/FTE and personnel cost/FTE.",
        ],
        0.7,
        1.5,
        12.0,
        5.5,
        18,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Data Pipeline", "Cleaning, standardization and metric assembly")
    add_bullets(
        slide,
        [
            "Ingestion: load 8 structured Excel sources across HR, training and finance.",
            "Normalization: trim column names, harmonize anonymized employee IDs, standardize dates and numeric fields.",
            "Entity layer: workforce snapshot becomes the reference grain for headcount, gender mix and country segmentation.",
            "Event layer: performance, training and absenteeism are transformed into employee-year or employee-event facts.",
            "Aggregation: KPI tables and charts are exported as CSV/PNG for notebook and presentation reuse.",
        ],
        0.7,
        1.6,
        5.3,
        5.2,
        18,
    )
    add_image(slide, FIG_DIR / "headcount_and_gender.png", 6.3, 1.6, 6.2)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Governance And Ethics", "Controls required before any human-capital AI layer")
    add_bullets(
        slide,
        [
            "GDPR: retain anonymized IDs only, keep purpose limitation explicit, and restrict outputs to aggregated insights for decision support.",
            "Data quality: track missing IDs, mixed date formats, duplicated employee-review records, and country-specific schema differences.",
            "Bias mitigation: monitor KPI differences by country and gender, and avoid using protected attributes as decision variables in downstream models.",
            "Human oversight: position the prototype as an augmentation tool for HRBP and managers, not an automated decision engine.",
        ],
        0.7,
        1.7,
        12.0,
        5.0,
        18,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "AI Prototype Direction", "Techniques aligned with the available data")
    add_bullets(
        slide,
        [
            "Descriptive layer: KPI monitoring with threshold-based alerts for review coverage, absence spikes and training gaps.",
            "Predictive layer: simple risk models for absenteeism concentration and missing review completion.",
            "Segmentation layer: clustering of employees or entities by training intensity, performance and well-being signals.",
            "Copilot layer: retrieval-augmented assistant that explains KPI movements and recommends follow-up actions with traceable evidence.",
        ],
        0.7,
        1.8,
        5.6,
        4.8,
        18,
    )
    add_image(slide, FIG_DIR / "training_vs_performance.png", 6.5, 1.7, 5.8)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Insight 1", "Scale stabilized after sharp growth, with a consistent inclusion baseline")
    add_image(slide, FIG_DIR / "headcount_and_gender.png", 0.7, 1.5, 7.0)
    latest = supporting["monthly_headcount"].iloc[-1]
    start_2024 = supporting["monthly_headcount"][supporting["monthly_headcount"]["PERIOD"].dt.year == 2024].iloc[0]
    add_bullets(
        slide,
        [
            f"Headcount moved from {int(start_2024['headcount'])} at the first 2024 snapshot to {int(latest['headcount'])} in the latest snapshot.",
            f"Female representation in the latest snapshot stands at {latest['female_share'] * 100:.1f}%, suggesting stability but still room for balance improvement.",
            "France and Luxembourg dominate the European footprint, so management practices in those geographies will drive most KPI shifts.",
        ],
        8.1,
        1.8,
        4.6,
        4.8,
        17,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Insight 2", "Capability signals are strong: large training volume and high learner satisfaction")
    add_image(slide, FIG_DIR / "training_hours.png", 0.7, 1.5, 5.8)
    add_image(slide, FIG_DIR / "training_effectiveness.png", 6.6, 1.5, 5.8)
    training_2025 = supporting["training_kpis"].query("year == 2025").iloc[0]
    rating_2025 = supporting["quick_scores"].query("year == 2025")
    impact_2025 = supporting["cold_impact"].query("year == 2025")
    rating_txt = f"{rating_2025['avg_training_rating'].iloc[0]:.2f}/5" if not rating_2025.empty else "n/a"
    impact_txt = f"{impact_2025['avg_positive_impact_rate'].iloc[0] * 100:.1f}%" if not impact_2025.empty else "n/a"
    add_bullets(
        slide,
        [
            f"Completed training reached {training_2025['training_hours']:.0f} hours in 2025 across {int(training_2025['trained_employees'])} employees.",
            f"Post-training satisfaction remains high at {rating_txt}, while delayed impact feedback shows {impact_txt} positive responses.",
            "This makes training one of the strongest candidates for the final AI use case because the signal is rich, repeated and linkable to employees.",
        ],
        0.8,
        5.7,
        12.0,
        1.0,
        16,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Insight 3", "Performance coverage is solid, and efficiency per FTE continues to improve")
    add_image(slide, FIG_DIR / "review_completion_by_country.png", 0.7, 1.5, 5.8)
    add_image(slide, FIG_DIR / "value_creation_efficiency.png", 6.6, 1.5, 5.8)
    review_overall = supporting["review_metrics"]["completion_rate"].mul(supporting["review_metrics"]["employees_with_reviews"]).sum() / supporting["review_metrics"]["employees_with_reviews"].sum()
    finance_2025 = supporting["finance"].query("year == 2025").iloc[0]
    add_bullets(
        slide,
        [
            f"Weighted 2024 review completion is {review_overall * 100:.1f}%, which is strong enough for KPI-driven managerial accountability.",
            f"2025 net banking income per FTE reaches {finance_2025['nbi_per_fte']:.0f}, while personnel cost per FTE is {finance_2025['personnel_cost_per_fte']:.0f}.",
            "The strategic opportunity is to connect capability and well-being signals to these business efficiency outcomes in the final prototype.",
        ],
        0.8,
        5.7,
        12.0,
        1.0,
        16,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Implementation Assets", "What is included in the technical submission")
    add_bullets(
        slide,
        [
            "Notebook: `deliverables/CACEIS_Human_Capital_Implementation.ipynb` with KPI logic and exported analysis tables.",
            "Pipeline script: `technical_implementation/generate_caceis_deliverables.py` to fully reproduce outputs.",
            "Reusable assets: CSV KPI tables and PNG charts in `deliverables/tables/` and `deliverables/figures/`.",
            "Next step for the final prototype: promote these tables into a dashboard and add interpretable risk scoring.",
        ],
        0.8,
        1.8,
        12.0,
        4.8,
        18,
    )

    prs.save(OUTPUT_DIR / "CACEIS_Human_Capital_Valuation_Presentation.pptx")


def make_summary_markdown(kpis: pd.DataFrame) -> None:
    lines = [
        "# CACEIS Human Capital Valuation Deliverables",
        "",
        "Generated files:",
        "- `deliverables/CACEIS_Human_Capital_Valuation_Presentation.pptx`",
        "- `deliverables/CACEIS_Human_Capital_Implementation.ipynb`",
        "- `deliverables/tables/*.csv`",
        "- `deliverables/figures/*.png`",
        "",
        "Headline KPIs:",
    ]
    for _, row in kpis.iterrows():
        lines.append(f"- {row['kpi']}: {row['value']} ({row['definition']})")
    (OUTPUT_DIR / "README.md").write_text("\n".join(lines), encoding="utf-8")


def main() -> None:
    ensure_dirs()
    data = load_sources()
    kpi_table, supporting = build_kpis(data)
    save_tables(kpi_table, supporting)
    save_figures(supporting)
    make_notebook()
    make_presentation(kpi_table, supporting)
    make_summary_markdown(kpi_table)


if __name__ == "__main__":
    main()
